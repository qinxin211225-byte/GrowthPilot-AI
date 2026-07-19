"""将 GrowthPilot 结果导出为企业常用办公文件。"""

from __future__ import annotations

from io import BytesIO
from typing import Iterable

import pandas as pd


class ExportUnavailableError(RuntimeError):
    pass


def _sections(markdown: str) -> list[tuple[str, str]]:
    title = "GrowthPilot AI OS 报告"
    content: list[str] = []
    sections: list[tuple[str, str]] = []
    for line in markdown.splitlines():
        if line.startswith("# ") or line.startswith("## "):
            if content:
                sections.append((title, "\n".join(content).strip()))
            title = line.lstrip("#").strip()
            content = []
        else:
            content.append(line)
    if content:
        sections.append((title, "\n".join(content).strip()))
    return sections or [(title, markdown)]


def _require(module_name: str, install_hint: str) -> None:
    try:
        __import__(module_name)
    except ImportError as exc:
        raise ExportUnavailableError(f"缺少导出依赖，请运行：pip install {install_hint}") from exc


def build_word_report(markdown: str, report_title: str = "GrowthPilot AI OS 报告") -> bytes:
    _require("docx", "python-docx")
    from docx import Document
    from docx.shared import Pt

    document = Document()
    normal_style = document.styles["Normal"]
    normal_style.font.name = "Microsoft YaHei"
    normal_style.font.size = Pt(10.5)
    document.add_heading(report_title, 0)
    for heading, body in _sections(markdown):
        document.add_heading(heading, level=1)
        for paragraph in body.split("\n\n"):
            document.add_paragraph(paragraph.strip())
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def build_pdf_report(markdown: str, report_title: str = "GrowthPilot AI OS 报告") -> bytes:
    _require("reportlab", "reportlab")
    from reportlab.lib.colors import HexColor
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    output = BytesIO()
    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "GrowthTitle", parent=styles["Title"], fontName="STSong-Light", fontSize=21,
        leading=29, textColor=HexColor("#312E81"), spaceAfter=13,
    )
    heading_style = ParagraphStyle(
        "GrowthHeading", parent=styles["Heading2"], fontName="STSong-Light", fontSize=14,
        leading=20, textColor=HexColor("#4F46E5"), spaceBefore=10, spaceAfter=6,
    )
    body_style = ParagraphStyle(
        "GrowthBody", parent=styles["BodyText"], fontName="STSong-Light", fontSize=9.6,
        leading=16, textColor=HexColor("#1F2937"),
    )
    story: list[object] = [Paragraph(report_title, title_style), Spacer(1, 3 * mm)]
    for heading, body in _sections(markdown):
        story.append(Paragraph(heading, heading_style))
        for paragraph in body.split("\n\n"):
            safe_text = paragraph.strip().replace("\n", "<br/>").replace("&", "&amp;")
            story.append(Paragraph(safe_text or " ", body_style))
            story.append(Spacer(1, 2 * mm))
    document = SimpleDocTemplate(output, pagesize=A4, rightMargin=18 * mm, leftMargin=18 * mm)
    document.build(story)
    return output.getvalue()


def build_ppt_report(markdown: str, report_title: str = "GrowthPilot AI OS 方案") -> bytes:
    _require("pptx", "python-pptx")
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    cover = presentation.slides.add_slide(presentation.slide_layouts[6])
    title_box = cover.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.6), Inches(1.4))
    title_run = title_box.text_frame.paragraphs[0].add_run()
    title_run.text = report_title
    title_run.font.name = "Microsoft YaHei"
    title_run.font.size = Pt(36)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(49, 46, 129)
    sub_box = cover.shapes.add_textbox(Inches(0.85), Inches(3.0), Inches(9.8), Inches(0.6))
    sub_run = sub_box.text_frame.paragraphs[0].add_run()
    sub_run.text = "GrowthPilot AI OS · 企业增长运营工作台"
    sub_run.font.name = "Microsoft YaHei"
    sub_run.font.size = Pt(18)
    sub_run.font.color.rgb = RGBColor(79, 70, 229)

    for heading, body in _sections(markdown):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        header = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(11.7), Inches(0.7))
        header_run = header.text_frame.paragraphs[0].add_run()
        header_run.text = heading
        header_run.font.name = "Microsoft YaHei"
        header_run.font.size = Pt(28)
        header_run.font.bold = True
        header_run.font.color.rgb = RGBColor(49, 46, 129)
        content_box = slide.shapes.add_textbox(Inches(0.82), Inches(1.55), Inches(11.55), Inches(5.25))
        frame = content_box.text_frame
        frame.word_wrap = True
        for index, paragraph in enumerate([item for item in body.split("\n\n") if item.strip()] or ["暂无内容"]):
            text_paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            text_paragraph.text = paragraph.replace("\n", " ")
            text_paragraph.font.name = "Microsoft YaHei"
            text_paragraph.font.size = Pt(17)
            text_paragraph.font.color.rgb = RGBColor(55, 65, 81)
            text_paragraph.space_after = Pt(12)
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def build_excel_report(
    markdown: str,
    dataframe: pd.DataFrame | None = None,
    report_title: str = "GrowthPilot AI OS 分析表",
) -> bytes:
    _require("openpyxl", "openpyxl")
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
    from openpyxl.utils import get_column_letter

    workbook = Workbook()
    summary = workbook.active
    summary.title = "分析报告"
    summary["A1"] = report_title
    summary["A1"].font = Font(name="Microsoft YaHei", size=16, bold=True, color="312E81")
    row = 3
    for heading, body in _sections(markdown):
        summary.cell(row=row, column=1, value=heading)
        summary.cell(row=row, column=1).font = Font(name="Microsoft YaHei", bold=True, color="FFFFFF")
        summary.cell(row=row, column=1).fill = PatternFill("solid", fgColor="6257D8")
        row += 1
        summary.cell(row=row, column=1, value=body)
        summary.cell(row=row, column=1).alignment = summary.cell(row=row, column=1).alignment.copy(wrap_text=True, vertical="top")
        row += 2
    summary.column_dimensions["A"].width = 105

    if dataframe is not None:
        raw = workbook.create_sheet("原始数据")
        for column_index, column in enumerate(dataframe.columns, start=1):
            cell = raw.cell(row=1, column=column_index, value=str(column))
            cell.font = Font(name="Microsoft YaHei", bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="312E81")
        for row_index, values in enumerate(dataframe.itertuples(index=False, name=None), start=2):
            for column_index, value in enumerate(values, start=1):
                raw.cell(row=row_index, column=column_index, value=value)
        for column_index in range(1, len(dataframe.columns) + 1):
            raw.column_dimensions[get_column_letter(column_index)].width = 18

    output = BytesIO()
    workbook.save(output)
    return output.getvalue()


def export_formats_available() -> dict[str, bool]:
    return {
        "word": _module_available("docx"),
        "pdf": _module_available("reportlab"),
        "ppt": _module_available("pptx"),
        "excel": _module_available("openpyxl"),
    }


def _module_available(module_name: str) -> bool:
    try:
        __import__(module_name)
        return True
    except ImportError:
        return False
